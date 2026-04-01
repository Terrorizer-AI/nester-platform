"""
Google Drive sync — fetches company docs, chunks them, embeds, and saves to SQLite.

Flow:
  1. List files in the configured Drive folder
  2. Compare with knowledge_files table (delta sync — only new/changed)
  3. Download + extract text from each changed file
  4. Chunk into ~500 token pieces with overlap
  5. Embed via OpenAI text-embedding-3-small
  6. Save chunks + embeddings to knowledge_chunks table
  7. Regenerate company master profile if any doc changed

Run standalone:  python knowledge/drive_sync.py
Called by:       start.sh on every startup (background)
                 api/routes/knowledge.py POST /knowledge/sync
"""

from __future__ import annotations

import io
import json
import logging
import os
import textwrap
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# ── Config ────────────────────────────────────────────────────────────────────

CHUNK_SIZE = 800        # characters per chunk (~200 tokens)
CHUNK_OVERLAP = 100     # overlap between chunks
EMBED_MODEL = "text-embedding-3-small"
EMBED_BATCH = 50        # embed up to 50 chunks per API call

SUPPORTED_MIME_TYPES = {
    "application/vnd.google-apps.document",
    "application/vnd.google-apps.presentation",
    "application/vnd.google-apps.spreadsheet",
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "text/plain",
    "text/markdown",
    "text/csv",
}

EXPORT_MIME = {
    "application/vnd.google-apps.document": "text/plain",
    "application/vnd.google-apps.presentation": "text/plain",
    "application/vnd.google-apps.spreadsheet": "text/csv",
}


# ── Drive client ──────────────────────────────────────────────────────────────

def _get_drive_service(credentials_path: str, token_path: str):
    """Build and return an authenticated Google Drive service."""
    from google.oauth2.credentials import Credentials
    from google.auth.transport.requests import Request
    from google_auth_oauthlib.flow import InstalledAppFlow
    from googleapiclient.discovery import build

    SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]
    creds = None

    token_file = Path(token_path)
    if token_file.exists():
        creds = Credentials.from_authorized_user_file(str(token_file), SCOPES)

    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(credentials_path, SCOPES)
            creds = flow.run_local_server(port=0)
        token_file.parent.mkdir(parents=True, exist_ok=True)
        token_file.write_text(creds.to_json())

    return build("drive", "v3", credentials=creds)


def _list_folder_files(service: Any, folder_id: str) -> list[dict]:
    """List all supported files in a Drive folder (non-recursive)."""
    results = []
    page_token = None

    mime_filter = " or ".join(f"mimeType='{m}'" for m in SUPPORTED_MIME_TYPES)
    query = f"'{folder_id}' in parents and trashed=false and ({mime_filter})"

    while True:
        resp = service.files().list(
            q=query,
            fields="nextPageToken, files(id, name, mimeType, modifiedTime)",
            pageToken=page_token,
        ).execute()

        results.extend(resp.get("files", []))
        page_token = resp.get("nextPageToken")
        if not page_token:
            break

    return results


def _download_file_text(service: Any, file_id: str, mime_type: str) -> str:
    """Download a Drive file and return its text content."""
    from googleapiclient.http import MediaIoBaseDownload

    export_mime = EXPORT_MIME.get(mime_type)

    if export_mime:
        # Google Workspace doc — export as plain text
        request = service.files().export_media(fileId=file_id, mimeType=export_mime)
    else:
        # Binary file (PDF, DOCX, etc.) — download raw
        request = service.files().get_media(fileId=file_id)

    buffer = io.BytesIO()
    downloader = MediaIoBaseDownload(buffer, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()

    raw_bytes = buffer.getvalue()

    if export_mime == "text/plain" or mime_type == "text/plain":
        return raw_bytes.decode("utf-8", errors="replace")

    if mime_type == "text/csv":
        return raw_bytes.decode("utf-8", errors="replace")

    if mime_type == "application/pdf":
        return _extract_pdf_text(raw_bytes)

    if mime_type in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ):
        return _extract_docx_text(raw_bytes)

    if mime_type in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ):
        return _extract_pptx_text(raw_bytes)

    # Fallback: try UTF-8 decode
    return raw_bytes.decode("utf-8", errors="replace")


def _extract_pdf_text(raw: bytes) -> str:
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(raw))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except ImportError:
        pass
    try:
        import pdfminer.high_level
        return pdfminer.high_level.extract_text(io.BytesIO(raw))
    except ImportError:
        logger.warning("[DriveSync] No PDF library — install pypdf or pdfminer.six")
        return ""


def _extract_docx_text(raw: bytes) -> str:
    try:
        import docx
        doc = docx.Document(io.BytesIO(raw))
        return "\n".join(p.text for p in doc.paragraphs)
    except ImportError:
        logger.warning("[DriveSync] python-docx not installed")
        return ""


def _extract_pptx_text(raw: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(raw))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except ImportError:
        logger.warning("[DriveSync] python-pptx not installed")
        return ""


# ── Chunking ──────────────────────────────────────────────────────────────────

def _chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Split text into overlapping chunks."""
    text = text.strip()
    if not text:
        return []

    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start += chunk_size - overlap

    return chunks


# ── Embeddings ────────────────────────────────────────────────────────────────

def _embed_chunks(chunks: list[str], api_key: str) -> list[list[float]]:
    """Embed a list of text chunks via OpenAI. Returns list of embedding vectors."""
    from openai import OpenAI
    client = OpenAI(api_key=api_key)

    all_embeddings: list[list[float]] = []
    for i in range(0, len(chunks), EMBED_BATCH):
        batch = chunks[i:i + EMBED_BATCH]
        response = client.embeddings.create(model=EMBED_MODEL, input=batch)
        all_embeddings.extend([item.embedding for item in response.data])

    return all_embeddings


# ── Main sync ─────────────────────────────────────────────────────────────────

def sync_drive(
    folder_id: str,
    credentials_path: str,
    token_path: str,
    openai_api_key: str,
    force_resync: bool = False,
) -> dict[str, Any]:
    """
    Sync Google Drive folder to local knowledge store.

    Returns summary: {synced, skipped, deleted, errors, profile_updated}
    """
    from memory.sqlite_ops import (
        get_knowledge_file, upsert_knowledge_chunk, upsert_knowledge_file,
        delete_knowledge_file, list_knowledge_files, is_sqlite_ready,
    )

    if not is_sqlite_ready():
        from memory.sqlite_ops import init_sqlite_ops
        init_sqlite_ops()

    summary = {"synced": 0, "skipped": 0, "deleted": 0, "errors": 0, "profile_updated": False}

    try:
        service = _get_drive_service(credentials_path, token_path)
    except Exception as e:
        logger.error("[DriveSync] Failed to connect to Google Drive: %s", e)
        summary["errors"] += 1
        return summary

    # List files in Drive folder
    try:
        drive_files = _list_folder_files(service, folder_id)
    except Exception as e:
        logger.error("[DriveSync] Failed to list Drive folder: %s", e)
        summary["errors"] += 1
        return summary

    drive_ids = {f["id"] for f in drive_files}

    # Delete removed files
    for stored in list_knowledge_files():
        if stored["file_id"] not in drive_ids:
            delete_knowledge_file(stored["file_id"])
            summary["deleted"] += 1
            logger.info("[DriveSync] Removed deleted file: %s", stored["file_name"])

    # Sync new/changed files
    any_changed = False
    for f in drive_files:
        file_id = f["id"]
        file_name = f["name"]
        mime_type = f["mimeType"]
        modified_time = f["modifiedTime"]

        existing = get_knowledge_file(file_id)
        if not force_resync and existing and existing["modified_time"] == modified_time:
            summary["skipped"] += 1
            continue

        logger.info("[DriveSync] Syncing: %s (%s)", file_name, mime_type)
        try:
            text = _download_file_text(service, file_id, mime_type)
            if not text.strip():
                logger.warning("[DriveSync] Empty content: %s", file_name)
                summary["skipped"] += 1
                continue

            chunks = _chunk_text(text)
            if not chunks:
                summary["skipped"] += 1
                continue

            embeddings = _embed_chunks(chunks, openai_api_key)

            # Delete old chunks first
            delete_knowledge_file(file_id)

            # Save new chunks
            for idx, (chunk, emb) in enumerate(zip(chunks, embeddings)):
                upsert_knowledge_chunk(file_id, file_name, idx, chunk, emb)

            upsert_knowledge_file(file_id, file_name, mime_type, modified_time, len(chunks))
            summary["synced"] += 1
            any_changed = True
            logger.info("[DriveSync] Synced %s → %d chunks", file_name, len(chunks))

        except Exception as e:
            logger.error("[DriveSync] Error syncing %s: %s", file_name, e)
            summary["errors"] += 1

    # Regenerate company master profile if anything changed
    if any_changed or force_resync:
        try:
            from knowledge.profile_builder import build_company_profile
            build_company_profile(openai_api_key)
            summary["profile_updated"] = True
        except Exception as e:
            logger.error("[DriveSync] Profile build failed: %s", e)

    logger.info("[DriveSync] Done — %s", summary)
    return summary


# ── File-picker based sync (uses OAuth access token, no credentials.json) ────

def sync_drive_files(
    file_ids: list[str],
    access_token: str,
    openai_api_key: str,
    force_resync: bool = False,
) -> dict[str, Any]:
    """
    Sync specific Drive files by ID using an OAuth access token.
    Used by the Google Picker flow — no credentials.json needed.

    Returns summary: {synced, skipped, errors, profile_updated}
    """
    from memory.sqlite_ops import (
        get_knowledge_file, upsert_knowledge_chunk, upsert_knowledge_file,
        delete_knowledge_file, is_sqlite_ready,
    )

    if not is_sqlite_ready():
        from memory.sqlite_ops import init_sqlite_ops
        init_sqlite_ops()

    summary = {"synced": 0, "skipped": 0, "errors": 0, "profile_updated": False}

    import requests

    any_changed = False
    for file_id in file_ids:
        # Get file metadata
        try:
            meta_resp = requests.get(
                f"https://www.googleapis.com/drive/v3/files/{file_id}",
                params={"fields": "id,name,mimeType,modifiedTime"},
                headers={"Authorization": f"Bearer {access_token}"},
            )
            meta_resp.raise_for_status()
            meta = meta_resp.json()
        except Exception as e:
            logger.error("[DriveSync] Failed to get metadata for %s: %s", file_id, e)
            summary["errors"] += 1
            continue

        file_name = meta["name"]
        mime_type = meta["mimeType"]
        modified_time = meta["modifiedTime"]

        if mime_type not in SUPPORTED_MIME_TYPES:
            logger.info("[DriveSync] Skipping unsupported type: %s (%s)", file_name, mime_type)
            summary["skipped"] += 1
            continue

        existing = get_knowledge_file(file_id)
        if not force_resync and existing and existing["modified_time"] == modified_time:
            summary["skipped"] += 1
            continue

        logger.info("[DriveSync] Syncing: %s (%s)", file_name, mime_type)
        try:
            export_mime = EXPORT_MIME.get(mime_type)
            if export_mime:
                dl_resp = requests.get(
                    f"https://www.googleapis.com/drive/v3/files/{file_id}/export",
                    params={"mimeType": export_mime},
                    headers={"Authorization": f"Bearer {access_token}"},
                )
            else:
                dl_resp = requests.get(
                    f"https://www.googleapis.com/drive/v3/files/{file_id}?alt=media",
                    headers={"Authorization": f"Bearer {access_token}"},
                )
            dl_resp.raise_for_status()
            raw_bytes = dl_resp.content

            if export_mime == "text/plain" or mime_type in ("text/plain", "text/markdown", "text/csv"):
                text = raw_bytes.decode("utf-8", errors="replace")
            elif mime_type == "application/pdf":
                text = _extract_pdf_text(raw_bytes)
            elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                text = _extract_docx_text(raw_bytes)
            elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                text = _extract_pptx_text(raw_bytes)
            else:
                text = raw_bytes.decode("utf-8", errors="replace")

            if not text.strip():
                logger.warning("[DriveSync] Empty content: %s", file_name)
                summary["skipped"] += 1
                continue

            chunks = _chunk_text(text)
            if not chunks:
                summary["skipped"] += 1
                continue

            embeddings = _embed_chunks(chunks, openai_api_key)
            delete_knowledge_file(file_id)

            for idx, (chunk, emb) in enumerate(zip(chunks, embeddings)):
                upsert_knowledge_chunk(file_id, file_name, idx, chunk, emb)

            upsert_knowledge_file(file_id, file_name, mime_type, modified_time, len(chunks))
            summary["synced"] += 1
            any_changed = True
            logger.info("[DriveSync] Synced %s → %d chunks", file_name, len(chunks))

        except Exception as e:
            logger.error("[DriveSync] Error syncing %s: %s", file_name, e)
            summary["errors"] += 1

    # Regenerate company master profile if anything changed
    if any_changed or force_resync:
        try:
            from knowledge.profile_builder import build_company_profile
            build_company_profile(openai_api_key)
            summary["profile_updated"] = True
        except Exception as e:
            logger.error("[DriveSync] Profile build failed: %s", e)

    logger.info("[DriveSync] Done — %s", summary)
    return summary


# ── Standalone entry point ────────────────────────────────────────────────────

if __name__ == "__main__":
    import sys
    from pathlib import Path

    _project_root = Path(__file__).parent.parent
    sys.path.insert(0, str(_project_root))

    from dotenv import load_dotenv
    load_dotenv(_project_root / ".env")

    from memory.sqlite_ops import init_sqlite_ops
    init_sqlite_ops()

    logging.basicConfig(level=logging.INFO, format="%(levelname)s %(message)s")

    folder_id = os.environ.get("GOOGLE_DRIVE_FOLDER_ID", "")
    creds_path = os.environ.get("GOOGLE_ACCOUNT_FILE", str(Path.home() / ".credentials/credentials.json"))
    token_path = os.environ.get("GOOGLE_TOKEN_FILE", str(Path.home() / ".credentials/gdrive_token.json"))
    api_key = os.environ.get("OPENAI_API_KEY", "")

    if not folder_id:
        print("ERROR: GOOGLE_DRIVE_FOLDER_ID not set in .env")
        sys.exit(1)
    if not api_key:
        print("ERROR: OPENAI_API_KEY not set in .env")
        sys.exit(1)

    force = "--force" in sys.argv
    result = sync_drive(folder_id, creds_path, token_path, api_key, force_resync=force)
    print(f"Sync complete: {result}")
